# -*- coding: utf-8 -*-
"""Build b11402037-b11402053-demo.pptx — C Programming term project demo deck."""
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY    = "1B2A41"   # dark slide bg
PANEL_D = "223349"   # panel on dark
CODE_BG = "15202F"   # code block bg
INK     = "22303F"   # body text on light
MUTED   = "5F6E7D"
TEAL    = "0E8C8B"
MINT    = "3FD0B6"
ICE     = "CADCFC"
LIGHT   = "F3F7F9"   # light panel
BORDER  = "D8E1E7"
WHITE   = "FFFFFF"
CODE_TX = "BFEFE3"
CODE_HI = "63E0C3"
CODE_CM = "7E93A8"
RED_SOFT= "C0564A"

ZH = "Microsoft JhengHei"
MONO = "Consolas"

PAGE_W, PAGE_H = 13.333, 7.5
MARGIN = 0.7

def C(hexs):
    return RGBColor.from_string(hexs)

prs = Presentation()
prs.slide_width  = Emu(int(PAGE_W * 914400))
prs.slide_height = Emu(int(PAGE_H * 914400))
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = C(bg)
    r.line.fill.background(); r.shadow.inherit = False
    return s

def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, In(x), In(y), In(w), In(h))
    if radius is not None:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = C(line); shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def arrow(slide, x, y, w, h, fill=TEAL):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, In(x), In(y), In(w), In(h))
    try:
        shp.adjustments[0] = 0.55; shp.adjustments[1] = 0.55
    except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def _style_run(run, text, size=13, color=INK, bold=False, italic=False,
               font=ZH, spc=None):
    run.text = text
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = C(color); f.name = font
    rPr = run._r.get_or_add_rPr()
    # east-asian font
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', ZH)
    if spc is not None:
        rPr.set('spc', str(spc))

def _bullet(p, color=TEAL, char="▪", indent=0.18):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(indent * 914400)))
    pPr.set('indent', str(-int(indent * 914400)))
    buClr = pPr.makeelement(qn('a:buClr'), {})
    srgb = pPr.makeelement(qn('a:srgbClr'), {'val': color})
    buClr.append(srgb); pPr.append(buClr)
    pPr.append(pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'}))
    pPr.append(pPr.makeelement(qn('a:buChar'), {'char': char}))

def text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts {runs:[(txt,opts)], align, space_after, line_spacing, bullet}"""
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get('align', PP_ALIGN.LEFT)
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        if para.get('space_before') is not None:
            p.space_before = Pt(para['space_before'])
        if para.get('line_spacing') is not None:
            p.line_spacing = para['line_spacing']
        for txt, opts in para['runs']:
            _style_run(p.add_run(), txt, **opts)
        if para.get('bullet'):
            _bullet(p, color=para.get('bullet_color', TEAL))
    return tb

def kicker_title(slide, kicker, title, color_k=TEAL, color_t=INK):
    text(slide, MARGIN, 0.42, 12.0, 0.3, [
        {'runs': [(kicker, dict(size=12, color=color_k, bold=True, font=MONO, spc=160))]},
    ])
    text(slide, MARGIN, 0.72, 12.0, 0.62, [
        {'runs': [(title, dict(size=27, color=color_t, bold=True))]},
    ])

PAGE_TOTAL = 9
def footer(slide, page_no, dark=False):
    c = ICE if dark else MUTED
    text(slide, MARGIN, 7.08, 8.0, 0.3, [
        {'runs': [("Clawer C Data Normalization Engine — b11402037 · b11402053",
                   dict(size=9, color=c, font=MONO))]},
    ])
    text(slide, 11.6, 7.08, 1.03, 0.3, [
        {'runs': [(f"{page_no} / {PAGE_TOTAL}", dict(size=9, color=c, font=MONO))],
         'align': PP_ALIGN.RIGHT},
    ])

def code_block(slide, x, y, w, h, lines, title=None, size=11.5):
    """lines: list of (text, color) tuples rendered in mono inside dark panel."""
    box(slide, x, y, w, h, fill=CODE_BG, radius=0.045)
    # three terminal dots
    for i, dc in enumerate(("E5705C", "E5C05C", "5CC98A")):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, In(x + 0.18 + i * 0.17), In(y + 0.14), In(0.09), In(0.09))
        d.fill.solid(); d.fill.fore_color.rgb = C(dc)
        d.line.fill.background(); d.shadow.inherit = False
    paras = []
    if title:
        paras.append({'runs': [(title, dict(size=10, color=CODE_CM, font=MONO))],
                      'space_after': 4})
    for ln, col in lines:
        paras.append({'runs': [(ln, dict(size=size, color=col, font=MONO))],
                      'line_spacing': 1.12})
    text(slide, x + 0.25, y + 0.36, w - 0.5, h - 0.5, paras)

def subhead(slide, x, y, label, color=INK, w=6.0):
    chip = box(slide, x, y + 0.05, 0.13, 0.13, fill=TEAL)
    text(slide, x + 0.24, y - 0.045, w, 0.34, [
        {'runs': [(label, dict(size=15.5, color=color, bold=True))]},
    ])

# =====================================================================
# Slide 1 — Cover
# =====================================================================
s = new_slide(bg=NAVY)
text(s, 0.9, 0.72, 11.5, 0.32, [
    {'runs': [("C PROGRAMMING · C 程式設計 — TERM PROJECT DEMO 專題展示",
               dict(size=12.5, color=MINT, bold=True, font=MONO, spc=140))]},
])
text(s, 0.9, 1.28, 11.6, 1.0, [
    {'runs': [("Clawer C Data Normalization Engine", dict(size=37, color=WHITE, bold=True))]},
])
text(s, 0.9, 2.08, 11.6, 0.55, [
    {'runs': [("C 語言大學排名資料正規化引擎", dict(size=21, color=ICE))]},
])
text(s, 0.9, 2.72, 11.6, 0.4, [
    {'runs': [("將異質爬蟲輸出（名稱／國家／排名／分數）轉換為可直接比較的標準化紀錄",
               dict(size=13, color=CODE_CM))]},
])

# member cards (left)
text(s, 0.9, 3.62, 5.0, 0.3, [
    {'runs': [("TEAM · 組員", dict(size=11, color=MINT, bold=True, font=MONO, spc=140))]},
])
members = [("b11402037", "高恩在"), ("b11402053", "林秉學")]
for i, (sid, name) in enumerate(members):
    yy = 4.02 + i * 0.86
    box(s, 0.9, yy, 4.7, 0.7, fill=PANEL_D, radius=0.10)
    text(s, 1.18, yy + 0.13, 2.3, 0.42, [
        {'runs': [(sid, dict(size=15, color=MINT, font=MONO, bold=True))]},
    ], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.45, yy + 0.13, 2.0, 0.42, [
        {'runs': [(name, dict(size=16, color=WHITE, bold=True))]},
    ], anchor=MSO_ANCHOR.MIDDLE)

# terminal panel (right)
code_block(s, 6.6, 3.62, 5.8, 2.5, [
    ("$ make run", CODE_HI),
    ("==== Clawer C Data Normalization Engine ====", CODE_TX),
    ("1. 載入   2. 正規化   3. 匯出", CODE_TX),
    ("", CODE_TX),
    ("$ make test_all", CODE_HI),
    ("Passed: 602 / 602 — All tests passed.", CODE_TX),
], size=11.5)

text(s, 0.9, 6.78, 11.5, 0.3, [
    {'runs': [("Demo：2026.06.18 – 06.20", dict(size=11, color=CODE_CM, font=MONO))]},
])

# =====================================================================
# Slide 2 — Introduction
# =====================================================================
s = new_slide()
kicker_title(s, "01 · INTRODUCTION", "簡介 — 背景、動機與研究問題")

subhead(s, MARGIN, 1.62, "背景與動機")
text(s, MARGIN, 2.0, 5.9, 0.75, [
    {'runs': [("上游爬蟲系統自 QS、THE、Shanghai Ranking 等多個排名來源抓取資料，"
               "輸出格式高度異質，無法直接進行跨來源比較與分析。",
               dict(size=12.5, color=INK))], 'line_spacing': 1.25},
])

dirty = [
    ("大學名稱", 'MIT · "University of   California… (UCB)"', "大小寫不一、縮寫、標點與空白差異"),
    ("國家名稱", "U.S.A. · UK · PR China · The Netherlands", "縮寫、官方長名、「The X」前綴變體"),
    ("排名格式", "53 · 101-150 · Top 100 · =201 · #10 · 201+", "純數字、區間、多種文字前綴"),
    ("分數欄位", '" 91.7 " · ".5" · "" · "N/A"', "空白、前導小數點、無效值"),
]
for i, (lab, sample, cap) in enumerate(dirty):
    yy = 2.82 + i * 0.98
    box(s, MARGIN, yy, 5.9, 0.86, fill=LIGHT, line=BORDER, line_w=0.75, radius=0.08)
    text(s, MARGIN + 0.22, yy + 0.10, 1.35, 0.3, [
        {'runs': [(lab, dict(size=12, color=TEAL, bold=True))]},
    ])
    text(s, MARGIN + 1.62, yy + 0.10, 4.1, 0.3, [
        {'runs': [(sample, dict(size=10.5, color=INK, font=MONO))]},
    ])
    text(s, MARGIN + 0.22, yy + 0.46, 5.5, 0.3, [
        {'runs': [(cap, dict(size=10.5, color=MUTED))]},
    ])

RX = 7.1; RW = 5.5
subhead(s, RX, 1.62, "研究問題")
box(s, RX, 2.0, RW, 1.46, fill="E8F4F2", radius=0.06)
box(s, RX, 2.0, 0.07, 1.46, fill=TEAL)
text(s, RX + 0.3, 2.18, RW - 0.55, 1.1, [
    {'runs': [("在不修改下游分析邏輯的前提下，如何將異質性爬蟲輸出轉換為可直接比較的數值紀錄？",
               dict(size=14, color=INK, bold=True))], 'line_spacing': 1.3},
])

subhead(s, RX, 3.82, "設計選擇 — 為何用 C？")
text(s, RX, 4.22, RW, 2.4, [
    {'runs': [("以 ", dict(size=12.5, color=INK)),
              ("C11（gcc -Wall -Wextra）", dict(size=12.5, color=TEAL, bold=True, font=MONO)),
              (" 實作核心引擎，而非 Python / pandas", dict(size=12.5, color=INK))],
     'bullet': True, 'space_after': 8, 'line_spacing': 1.2},
    {'runs': [("大量 CSV 批次處理時維持低記憶體佔用與可預期的執行時間",
               dict(size=12.5, color=INK))], 'bullet': True, 'space_after': 8, 'line_spacing': 1.2},
    {'runs': [("無 Python 執行期依賴，便於與上游爬蟲管線部署整合",
               dict(size=12.5, color=INK))], 'bullet': True, 'space_after': 8, 'line_spacing': 1.2},
    {'runs': [("單機、檔案導向：無資料庫伺服器、無 Web API、無網路依賴",
               dict(size=12.5, color=INK))], 'bullet': True, 'line_spacing': 1.2},
])
footer(s, 2)

# =====================================================================
# Slide 3 — Database Design (1): schema & data flow
# =====================================================================
s = new_slide()
kicker_title(s, "02 · DATABASE DESIGN", "資料庫設計 — CSV Schema 與資料流")
text(s, MARGIN, 1.5, 11.9, 0.36, [
    {'runs': [("採", dict(size=12.5, color=INK)),
              ("檔案導向", dict(size=12.5, color=TEAL, bold=True)),
              ("設計：以 CSV 作為持久層，不依賴外部資料庫伺服器，降低部署複雜度並貼合爬蟲批次輸出情境。",
               dict(size=12.5, color=INK))]},
])

text(s, MARGIN, 2.02, 5.4, 0.3, [
    {'runs': [("輸入 — 爬蟲原始 CSV（11 欄）", dict(size=12.5, color=INK, bold=True))]},
])
code_block(s, MARGIN, 2.34, 5.4, 1.42, [
    ("QS Rank,University,Country,GMAT,GRE,", CODE_TX),
    ("GPA,IELTS,TOEFL,Duolingo,", CODE_TX),
    ("Overall Score,URL", CODE_TX),
], size=11.5)

arrow(s, 6.28, 2.83, 0.66, 0.42)
text(s, 6.02, 3.32, 1.2, 0.55, [
    {'runs': [("正規化", dict(size=10, color=TEAL, bold=True))], 'align': PP_ALIGN.CENTER},
])

text(s, 7.1, 2.02, 5.5, 0.3, [
    {'runs': [("輸出 — 標準化 CSV（5 欄）", dict(size=12.5, color=INK, bold=True))]},
])
code_block(s, 7.1, 2.34, 5.53, 1.42, [
    ("University,Country,", CODE_TX),
    ("Rank Min,Rank Max,", CODE_TX),
    ("Overall Score", CODE_TX),
], size=11.5)

text(s, MARGIN, 4.0, 5.0, 0.3, [
    {'runs': [("實際資料範例（前 → 後）", dict(size=12.5, color=INK, bold=True))]},
])
code_block(s, MARGIN, 4.32, 11.93, 1.18, [
    ('raw  | University of   California Berkeley (UCB)  , usa ,12 , 14 ,  91.7', RED_SOFT),
    ('norm | University of California Berkeley (Ucb),United States,12,14,91.70', CODE_HI),
], size=11.5)

text(s, MARGIN, 5.72, 11.93, 1.1, [
    {'runs': [("Reader 防衛機制：header 欄位數與欄名驗證（不符即拒絕載入）、RFC 4180 quoted CSV 解析、UTF-8 BOM 清理、不合法列略過",
               dict(size=12, color=INK))], 'bullet': True, 'space_after': 7},
    {'runs': [("Writer 輸出保障：必要時自動 quoting / escaping，數值欄固定 %.2f 格式",
               dict(size=12, color=INK))], 'bullet': True},
])
footer(s, 3)

# =====================================================================
# Slide 4 — Database Design (2): core data structures
# =====================================================================
s = new_slide()
kicker_title(s, "02 · DATABASE DESIGN", "資料庫設計 — 核心資料結構與對照表")

code_block(s, MARGIN, 1.62, 5.95, 4.2, [
    ("/* include/record.h */", CODE_CM),
    ("typedef struct {", CODE_TX),
    ("    /* Raw input fields */", CODE_CM),
    ("    char raw_name[150];", CODE_TX),
    ("    char raw_country[80];", CODE_TX),
    ("    char raw_rank[50];", CODE_TX),
    ("    char raw_score[50];", CODE_TX),
    ("    /* Normalized fields */", CODE_CM),
    ("    char normalized_name[150];", CODE_TX),
    ("    char normalized_country[80];", CODE_TX),
    ("    int    rank_min, rank_max;", CODE_HI),
    ("    double score;", CODE_HI),
    ("} UniversityRecord;", CODE_TX),
], size=11.5)
text(s, MARGIN, 5.98, 5.95, 0.5, [
    {'runs': [("raw 與 normalized 欄位並存於同一筆紀錄，便於前後對照、預覽與除錯",
               dict(size=11, color=MUTED))], 'line_spacing': 1.2},
])

cards = [
    ("國家別名對照表 mapping_table[]",
     "靜態 key → value 查表：先 normalize_basic（去標點、轉小寫、壓縮空白）再比對；涵蓋縮寫（U.S.A.）、官方長名（Federal Republic of Germany）與「The X」變體，38+ 條目（Session 6 擴充）。"),
    ("排名區間表示 rank_min / rank_max",
     "以兩個整數欄位統一表達單值與區間：53 → [53,53]；101-150 → [101,150]；Top 100 → [1,100]，下游可直接做數值比較。"),
    ("哨兵值設計 score = -1",
     "無法解析的分數標記為 -1，由 writer 層轉為空欄位輸出，避免髒資料混入下游統計。"),
]
for i, (t, b) in enumerate(cards):
    yy = 1.62 + i * 1.78
    box(s, 7.05, yy, 5.58, 1.62, fill=LIGHT, line=BORDER, line_w=0.75, radius=0.06)
    box(s, 7.05, yy, 0.07, 1.62, fill=TEAL)
    text(s, 7.32, yy + 0.14, 5.1, 0.3, [
        {'runs': [(t, dict(size=13, color=INK, bold=True))]},
    ])
    text(s, 7.32, yy + 0.5, 5.1, 1.05, [
        {'runs': [(b, dict(size=10.8, color=MUTED))], 'line_spacing': 1.18},
    ])
footer(s, 4)

# =====================================================================
# Slide 5 — Implementation (1): architecture
# =====================================================================
s = new_slide()
kicker_title(s, "03 · IMPLEMENTATION", "實作 — 系統架構與模組分工")

flow_y, small_h, pipe_h = 1.66, 1.0, 2.62
pipe_y = flow_y
small_y = flow_y + (pipe_h - small_h) / 2
xx = MARGIN
def flow_box(label_lines, w, small=True, mono_first=True):
    global xx
    yy = small_y if small else pipe_y
    hh = small_h if small else pipe_h
    box(s, xx, yy, w, hh, fill=LIGHT, line=BORDER, line_w=1.0, radius=0.07)
    paras = []
    for j, ln in enumerate(label_lines):
        is_mono = mono_first and j == 0
        paras.append({'runs': [(ln, dict(size=10.5 if is_mono else 9.5,
                                         color=INK if is_mono else MUTED,
                                         bold=is_mono, font=MONO if is_mono else ZH))],
                      'align': PP_ALIGN.CENTER, 'line_spacing': 1.05})
    text(s, xx + 0.06, yy, w - 0.12, hh, paras, anchor=MSO_ANCHOR.MIDDLE)
    xx += w

def flow_arrow():
    global xx
    arrow(s, xx + 0.02, small_y + small_h / 2 - 0.14, 0.24, 0.28)
    xx += 0.28

flow_box(["raw_*.csv", "爬蟲輸出"], 1.2)
flow_arrow()
flow_box(["csv_reader.c", "header 驗證 / BOM"], 1.8)
flow_arrow()
flow_box(["UniversityRecord[]", "記憶體陣列"], 1.9)
flow_arrow()
# pipeline tall box
pw = 2.3
box(s, xx, pipe_y, pw, pipe_h, fill="E8F4F2", line=TEAL, line_w=1.2, radius=0.05)
text(s, xx, pipe_y + 0.12, pw, 0.3, [
    {'runs': [("normalizer.c", dict(size=11, color=TEAL, bold=True, font=MONO))],
     'align': PP_ALIGN.CENTER},
])
subs = ["name_normalizer.c", "country_normalizer.c", "rank_parser.c", "score_parser.c"]
for j, sub in enumerate(subs):
    syy = pipe_y + 0.5 + j * 0.51
    box(s, xx + 0.15, syy, pw - 0.3, 0.42, fill=WHITE, line=BORDER, line_w=0.75, radius=0.10)
    text(s, xx + 0.15, syy + 0.08, pw - 0.3, 0.3, [
        {'runs': [(sub, dict(size=9, color=INK, font=MONO))], 'align': PP_ALIGN.CENTER},
    ])
xx += pw
flow_arrow()
flow_box(["csv_writer.c", "quoting / escaping"], 1.8)
flow_arrow()
flow_box(["normalized_*.csv", "標準化輸出"], 1.55)

# bottom panels
by = 4.78
box(s, MARGIN, by, 5.9, 1.78, fill=CODE_BG, radius=0.05)
text(s, MARGIN + 0.25, by + 0.18, 5.4, 0.3, [
    {'runs': [("互動式 CLI — src/main.c", dict(size=11.5, color=CODE_HI, bold=True, font=MONO))]},
])
text(s, MARGIN + 0.25, by + 0.56, 5.4, 1.0, [
    {'runs': [("1 載入＋顯示原始   2 正規化＋顯示結果", dict(size=11.5, color=CODE_TX, font=MONO))],
     'line_spacing': 1.35},
    {'runs': [("3 匯出 CSV＋自動離開   0 離開", dict(size=11.5, color=CODE_TX, font=MONO))],
     'line_spacing': 1.35},
    {'runs': [("同步複合動作設計；狀態防呆擋下未就緒操作", dict(size=10.5, color=CODE_CM))],
     'line_spacing': 1.35},
])

box(s, 7.05, by, 5.58, 1.78, fill=LIGHT, line=BORDER, line_w=0.75, radius=0.05)
text(s, 7.3, by + 0.18, 5.1, 0.3, [
    {'runs': [("共用工具與建置", dict(size=12.5, color=INK, bold=True))]},
])
text(s, 7.3, by + 0.56, 5.15, 1.1, [
    {'runs': [("utils.c", dict(size=11, color=TEAL, bold=True, font=MONO)),
              ("：trim、collapse spaces、title case、縮寫（MIT / ETH / KAUST）偵測",
               dict(size=11, color=INK))], 'line_spacing': 1.25, 'space_after': 5},
    {'runs': [("Makefile", dict(size=11, color=TEAL, bold=True, font=MONO)),
              ("：make run 互動執行 · make test 單元測試 · 測試套件一鍵建置",
               dict(size=11, color=INK))], 'line_spacing': 1.25},
])
footer(s, 5)

# =====================================================================
# Slide 6 — Implementation (2): normalization rules table
# =====================================================================
s = new_slide()
kicker_title(s, "03 · IMPLEMENTATION", "實作 — 正規化規則與實例")

rows = [
    ("模組", "原始輸入（raw）", "正規化輸出（normalized）"),
    ("名稱", 'University of   California Berkeley (UCB)', "University of California Berkeley (Ucb)"),
    ("名稱", "KAUST · hkust · eth", "KAUST · HKUST · ETH（知名縮寫保留全大寫）"),
    ("國家", "usa · U.S.A. · UK", "United States · United Kingdom"),
    ("國家", "The Netherlands · PR China", "Netherlands · China (Mainland)"),
    ("排名", "Top 100 · #10 · =201", "[1,100] · [10,10] · [201,201]"),
    ("排名", "101-150 · 150-101", "[101,150]（反轉區間自動交換）"),
    ("分數", '" 91.7 " · ".5"', "91.70 · 0.50（前導小數點修正）"),
    ("分數", '"" · "N/A" · "abc"', "-1（哨兵值 → 輸出層轉空欄）"),
]
n_rows, n_cols = len(rows), 3
tbl_shape = s.shapes.add_table(n_rows, n_cols, In(MARGIN), In(1.6), In(11.93), In(4.86))
tbl = tbl_shape.table
tbl.first_row = False
tbl.horz_banding = False
tbl.columns[0].width = In(1.3)
tbl.columns[1].width = In(5.0)
tbl.columns[2].width = In(5.63)
for r in range(n_rows):
    tbl.rows[r].height = In(0.54)
    for c_i in range(n_cols):
        cell = tbl.cell(r, c_i)
        cell.margin_left = In(0.14); cell.margin_right = In(0.08)
        cell.margin_top = In(0.03); cell.margin_bottom = In(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = C(NAVY)
            _style_run(p.add_run(), rows[r][c_i], size=12, color=WHITE, bold=True)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C(WHITE if r % 2 == 1 else "F0F5F7")
            if c_i == 0:
                _style_run(p.add_run(), rows[r][c_i], size=11.5, color=TEAL, bold=True)
            elif c_i == 1:
                _style_run(p.add_run(), rows[r][c_i], size=10.5, color=RED_SOFT, font=MONO)
            else:
                _style_run(p.add_run(), rows[r][c_i], size=10.5, color=INK, font=MONO)

text(s, MARGIN, 6.62, 11.93, 0.35, [
    {'runs': [("規則涵蓋 6 種排名格式、大小寫不敏感比對（\"TOP 100\" ≡ \"Top 100\"），並拒絕負數排名、TOP 0 與整數溢位（strtol + 上限防護）。",
               dict(size=11, color=MUTED))]},
])
footer(s, 6)

# =====================================================================
# Slide 7 — Implementation (3): testing & verification
# =====================================================================
s = new_slide()
kicker_title(s, "03 · IMPLEMENTATION", "實作 — 弱點導向測試與驗證")

subhead(s, MARGIN, 1.58, "測試方法論（每輪 Session）")
steps = [
    ("靜態分析", "通讀模組程式碼，標記可疑邏輯"),
    ("探針程式", "撰寫 probe 重現可疑行為"),
    ("修正", "最小修改修復根因"),
    ("回歸驗證", "新增 regression 套件全量重跑"),
]
for i, (t, d) in enumerate(steps):
    yy = 2.02 + i * 0.78
    box(s, MARGIN, yy, 0.52, 0.52, fill=NAVY, radius=0.14)
    text(s, MARGIN, yy + 0.075, 0.52, 0.36, [
        {'runs': [(str(i + 1), dict(size=15, color=MINT, bold=True, font=MONO))],
         'align': PP_ALIGN.CENTER},
    ])
    text(s, MARGIN + 0.72, yy - 0.02, 4.9, 0.3, [
        {'runs': [(t, dict(size=13.5, color=INK, bold=True))]},
    ])
    text(s, MARGIN + 0.72, yy + 0.28, 4.9, 0.3, [
        {'runs': [(d, dict(size=10.5, color=MUTED))]},
    ])

box(s, MARGIN, 5.32, 5.9, 1.28, fill="E8F4F2", radius=0.06)
box(s, MARGIN, 5.32, 0.07, 1.28, fill=TEAL)
text(s, MARGIN + 0.28, 5.46, 5.45, 1.0, [
    {'runs': [("累計 6 輪 Session、17 項修正（A–T）", dict(size=12.5, color=INK, bold=True))],
     'space_after': 4},
    {'runs': [("例：Fix S — parse_score(\".5\") 誤判為 5.0 → 修正為 0.5",
               dict(size=10.5, color=MUTED, font=MONO))]},
])

RX = 7.05; RW = 5.58
text(s, RX, 1.52, RW, 0.85, [
    {'runs': [("602 / 602", dict(size=48, color=TEAL, bold=True, font=MONO))]},
])
text(s, RX, 2.42, RW, 0.3, [
    {'runs': [("九個測試套件全數通過（make test_all）", dict(size=12.5, color=INK, bold=True))]},
])

suites = [
    ("test_normalizer", "17 / 17"),
    ("test_extreme", "139 / 139"),
    ("test_scale", "31 / 31"),
    ("test_weakness", "43 / 43"),
    ("test_regression2 – 6", "372 / 372"),
]
ty = 2.92
for i, (name, score) in enumerate(suites):
    yy = ty + i * 0.6
    box(s, RX, yy, RW, 0.5, fill=LIGHT if i % 2 == 0 else WHITE,
        line=BORDER, line_w=0.5)
    text(s, RX + 0.2, yy + 0.1, 3.6, 0.3, [
        {'runs': [(name, dict(size=11.5, color=INK, font=MONO))]},
    ])
    text(s, RX + RW - 1.85, yy + 0.1, 1.65, 0.3, [
        {'runs': [(score, dict(size=11.5, color=TEAL, bold=True, font=MONO))],
         'align': PP_ALIGN.RIGHT},
    ])
box(s, RX, ty + 5 * 0.6, RW, 0.56, fill=NAVY)
text(s, RX + 0.2, ty + 5 * 0.6 + 0.12, 3.6, 0.3, [
    {'runs': [("合計", dict(size=12, color=WHITE, bold=True))]},
])
text(s, RX + RW - 1.85, ty + 5 * 0.6 + 0.12, 1.65, 0.3, [
    {'runs': [("602 / 602", dict(size=12, color=MINT, bold=True, font=MONO))],
     'align': PP_ALIGN.RIGHT},
])
footer(s, 7)

# =====================================================================
# Slide 8 — Conclusion
# =====================================================================
s = new_slide()
kicker_title(s, "04 · CONCLUSION", "結論 — 成果、限制與未來工作")

cols = [
    ("主要成果", TEAL, [
        "完成 C11 正規化引擎：reader → 四大解析模組 → writer 完整 pipeline",
        "602 項單元／回歸測試全數通過，涵蓋極端值、規模與弱點情境",
        "6 輪弱點導向修正（A–T 共 17 項）建立回歸防護網",
    ]),
    ("已知限制", RED_SOFT, [
        "ASCII-only：É / ä / ã 等非 ASCII 字元直接傳遞，不做 Unicode 正規化",
        "requirement_parser 尚未實作（GMAT / GRE / GPA / IELTS / TOEFL / Duolingo）",
        "%.2f 捨入：99.999 輸出為 100.00（超出滿分）",
    ]),
    ("未來工作", NAVY, [
        "引入 utf8proc / ICU 處理非 ASCII 字元正規化",
        "國家別名表外部化：從 mapping 檔載入，取代硬編碼",
        "CI 整合（爬蟲 → 引擎 → 資料庫）並提供 libnormalizer 函式庫介面",
    ]),
]
cw = 3.84; gap = 0.2
for i, (title, color, items) in enumerate(cols):
    xx2 = MARGIN + i * (cw + gap)
    box(s, xx2, 1.62, cw, 4.3, fill=LIGHT, line=BORDER, line_w=0.75, radius=0.05)
    box(s, xx2, 1.62, cw, 0.62, fill=color)
    text(s, xx2 + 0.25, 1.76, cw - 0.5, 0.34, [
        {'runs': [(title, dict(size=14.5, color=WHITE, bold=True))]},
    ])
    paras = []
    for it in items:
        paras.append({'runs': [(it, dict(size=11.5, color=INK))],
                      'bullet': True, 'bullet_color': color,
                      'space_after': 10, 'line_spacing': 1.22})
    text(s, xx2 + 0.28, 2.5, cw - 0.56, 3.3, paras)

box(s, MARGIN, 6.18, 11.93, 0.62, fill=NAVY, radius=0.08)
text(s, MARGIN, 6.33, 11.93, 0.34, [
    {'runs': [("以小而可驗證的 C 核心，換取資料管線中可預期的正規化行為。",
               dict(size=14, color=MINT, bold=True))], 'align': PP_ALIGN.CENTER},
])
footer(s, 8)

# =====================================================================
# Slide 9 — References
# =====================================================================
s = new_slide(bg=NAVY)
text(s, 0.9, 0.5, 11.5, 0.3, [
    {'runs': [("REFERENCES", dict(size=12, color=MINT, bold=True, font=MONO, spc=160))]},
])
text(s, 0.9, 0.8, 11.5, 0.6, [
    {'runs': [("參考文獻", dict(size=27, color=WHITE, bold=True))]},
])

groups = [
    ("排名資料來源", [
        "QS World University Rankings — topuniversities.com",
        "Times Higher Education (THE) World University Rankings",
        "ShanghaiRanking / Academic Ranking of World Universities (ARWU)",
    ]),
    ("技術標準", [
        "ISO 3166-1 — 國家代碼標準（國家別名映射依據）",
        "RFC 4180 — CSV 格式標準（quoting / escaping 規範）",
        "Unicode UAX #15 — Unicode Normalization Forms（NFC / NFD）",
    ]),
    ("語言與工具", [
        "ISO/IEC 9899:2011 — C11 Language Standard",
        "GCC Compiler Documentation（-Wall -Wextra）",
        "GNU Make Manual",
    ]),
    ("專案內部文件", [
        "README.md — 使用說明與驗證流程",
        "c_engine/docs/architecture.md — 系統架構與資料流",
        "test-report.md — Session 1–6 累計測試報告",
    ]),
]
positions = [(0.9, 1.78), (7.0, 1.78), (0.9, 4.18), (7.0, 4.18)]
for (gx, gy), (gt, items) in zip(positions, groups):
    box(s, gx, gy + 0.05, 0.13, 0.13, fill=MINT)
    text(s, gx + 0.26, gy - 0.04, 5.3, 0.32, [
        {'runs': [(gt, dict(size=14.5, color=WHITE, bold=True))]},
    ])
    paras = []
    for it in items:
        paras.append({'runs': [(it, dict(size=11, color=ICE))],
                      'bullet': True, 'bullet_color': MINT,
                      'space_after': 7, 'line_spacing': 1.15})
    text(s, gx + 0.26, gy + 0.42, 5.5, 1.7, paras)

box(s, 0.9, 6.52, 11.53, 0.0)
text(s, 0.9, 6.6, 11.5, 0.4, [
    {'runs': [("Thank You — 謝謝聆聽   |   Demo：2026.06.18 – 06.20   |   b11402037 高恩在 · b11402053 林秉學",
               dict(size=12, color=CODE_CM, font=MONO))]},
])
footer(s, 9, dark=True)

OUT = r"C:\Users\User\OneDrive\Desktop\Clawer-C-Data-Normalization-Engine\b11402037-b11402053-demo.pptx"
prs.save(OUT)
print("saved:", OUT)
